#!/usr/bin/env python3
"""
letras_to_pptx.py

Gera uma apresentação PPTX (9:16, retrato) a partir de um arquivo .txt
onde cada "### Título" marca o início de um slide, e o texto abaixo dele
(até o próximo ###) é o conteúdo (letra da música).

Regras de parsing/layout:
- Cada "###" vira um novo slide (título do slide).
- Dentro de cada slide, o conteúdo é dividido em PARÁGRAFOS por linha em
  branco (uma estrofe = um parágrafo). As quebras de linha DENTRO de um
  parágrafo são preservadas (não são "reflowed").
- O layout tenta usar a fonte grande (perto do tamanho do título). Se o
  conteúdo não couber, ele:
    1) tenta reduzir a fonte da letra (dentro de uma faixa definida em
       FONT_SIZES_LYRIC);
    2) se mesmo assim não couber TODOS os parágrafos, divide o conteúdo
       em vários slides, sempre cortando entre parágrafos (nunca no meio
       de um parágrafo) e nomeia como "Título - Parte 1", "Parte 2"...

Uso:
    python3 letras_to_pptx.py entrada.txt saida.pptx

Dependências:
    pip install python-pptx
"""

import argparse
import re
import sys
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURAÇÃO — mexa aqui pra ajustar tema, fontes e margens
# ============================================================

# Slide 9:16 (retrato). 7.5" x 13.333" mantém a proporção do 16:9 padrão,
# só girada.
SLIDE_WIDTH_IN = 7.5
SLIDE_HEIGHT_IN = 13.333

# Margens da área de conteúdo
MARGIN_X_IN = 0.5
MARGIN_TOP_IN = 0.6
MARGIN_BOTTOM_IN = 0.5

# Título
TITLE_HEIGHT_IN = 1.6
TITLE_FONT_SIZE = 60
TITLE_FONT_NAME = "Anton"          # fonte "poster", trocar se não tiver no PC
TITLE_COLOR = RGBColor(0xE8, 0x1F, 0x1F)   # vermelho rock
TITLE_ALIGN = PP_ALIGN.LEFT

# Letra da música — tenta do maior pro menor até caber
FONT_SIZES_LYRIC = [50, 44, 40, 36, 32, 28, 24, 20]
LYRIC_FONT_NAME = "Anton"
LYRIC_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
LYRIC_ALIGN = PP_ALIGN.LEFT
LINE_SPACING_FACTOR = 1.12          # altura da linha = font_size * fator
PARAGRAPH_GAP_PT = 14               # espaço extra entre estrofes

# Fundo
BG_COLOR = RGBColor(0x12, 0x12, 0x12)   # quase preto

# Heurística de largura média de caractere (fonte bold condensada tipo
# Anton/Oswald). Ajuste se o texto estourar ou sobrar muito espaço na
# hora de abrir o PPTX de verdade.
CHAR_WIDTH_RATIO = 0.56   # largura média do caractere = font_size(pt) * ratio, em pt

# Margem de segurança: encolhe a caixa disponível antes de medir, pra
# preferir dividir slide a estourar texto.
SAFETY_MARGIN = 0.90

# ============================================================
# PARSING DO .TXT
# ============================================================


@dataclass
class Slide:
    title: str
    paragraphs: list  # lista de parágrafos; cada parágrafo é lista de linhas (str)


def parse_input(text: str) -> list:
    """Divide o texto em slides por '### título' e cada slide em parágrafos
    (separados por linha em branco)."""
    # Garante que o arquivo comece com um ###, senão ignora preâmbulo
    blocks = re.split(r"(?m)^###\s*", text)
    slides = []
    for block in blocks:
        if not block.strip():
            continue
        lines = block.splitlines()
        title = lines[0].strip()
        body = "\n".join(lines[1:]).strip("\n")

        # separa em parágrafos por linha(s) em branco
        raw_paragraphs = re.split(r"\n\s*\n", body)
        paragraphs = []
        for rp in raw_paragraphs:
            rp = rp.strip("\n")
            if not rp.strip():
                continue
            paragraphs.append([ln for ln in rp.split("\n")])

        if paragraphs:
            slides.append(Slide(title=title, paragraphs=paragraphs))
    return slides


# ============================================================
# MEDIÇÃO DE TEXTO (heurística, sem dependências externas)
# ============================================================


def wrap_line_by_width(line: str, font_size_pt: float, box_width_in: float) -> list:
    """Quebra uma linha em sub-linhas que cabem em box_width_in, usando a
    largura média de caractere estimada. Palavras não são cortadas."""
    if line.strip() == "":
        return [""]

    avg_char_width_pt = font_size_pt * CHAR_WIDTH_RATIO
    avg_char_width_in = avg_char_width_pt / 72.0
    if avg_char_width_in <= 0:
        return [line]

    chars_per_line = max(1, int(box_width_in / avg_char_width_in))

    words = line.split(" ")
    wrapped = []
    current = ""
    for word in words:
        candidate = word if not current else current + " " + word
        if len(candidate) <= chars_per_line:
            current = candidate
        else:
            if current:
                wrapped.append(current)
            # palavra sozinha maior que a linha: força quebra bruta
            while len(word) > chars_per_line:
                wrapped.append(word[:chars_per_line])
                word = word[chars_per_line:]
            current = word
    if current:
        wrapped.append(current)
    return wrapped or [""]


def paragraph_height_in(paragraph_lines: list, font_size_pt: float, box_width_in: float) -> float:
    line_height_in = (font_size_pt * LINE_SPACING_FACTOR) / 72.0
    total_lines = 0
    for line in paragraph_lines:
        total_lines += len(wrap_line_by_width(line, font_size_pt, box_width_in))
    return total_lines * line_height_in


def paragraphs_block_height_in(paragraphs: list, font_size_pt: float, box_width_in: float) -> float:
    gap_in = PARAGRAPH_GAP_PT / 72.0
    total = 0.0
    for i, para in enumerate(paragraphs):
        total += paragraph_height_in(para, font_size_pt, box_width_in)
        if i < len(paragraphs) - 1:
            total += gap_in
    return total


# ============================================================
# QUEBRA EM SLIDES (por parágrafo, nunca no meio de um parágrafo)
# ============================================================


def plan_slides_for_song(paragraphs: list, box_width_in: float, box_height_in: float):
    """Retorna lista de tuplas (lista_de_paragrafos, font_size) — uma
    entrada por slide necessário para caber a música inteira."""
    avail_w = box_width_in * SAFETY_MARGIN
    avail_h = box_height_in * SAFETY_MARGIN

    plan = []
    remaining = list(paragraphs)

    while remaining:
        chosen_font = None
        chosen_count = 0

        for font_size in FONT_SIZES_LYRIC:
            count = 0
            for n in range(1, len(remaining) + 1):
                h = paragraphs_block_height_in(remaining[:n], font_size, avail_w)
                if h <= avail_h:
                    count = n
                else:
                    break
            if count > 0:
                chosen_font = font_size
                chosen_count = count
                break  # maior fonte que já encaixa pelo menos 1 parágrafo

        if chosen_font is None:
            # nem o parágrafo sozinho cabe na menor fonte -> manda mesmo
            # assim (evita loop infinito); vai estourar um pouco.
            chosen_font = FONT_SIZES_LYRIC[-1]
            chosen_count = 1

        plan.append((remaining[:chosen_count], chosen_font))
        remaining = remaining[chosen_count:]

    return plan


# ============================================================
# GERAÇÃO DO PPTX
# ============================================================


def add_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    bg.shadow.inherit = False
    # manda pro fundo
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_title(slide, prs, title_text):
    left = Inches(MARGIN_X_IN)
    top = Inches(MARGIN_TOP_IN)
    width = Inches(SLIDE_WIDTH_IN - 2 * MARGIN_X_IN)
    height = Inches(TITLE_HEIGHT_IN)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = TITLE_ALIGN
    run = p.add_run()
    run.text = title_text.upper()
    run.font.size = Pt(TITLE_FONT_SIZE)
    run.font.bold = True
    run.font.name = TITLE_FONT_NAME
    run.font.color.rgb = TITLE_COLOR


def add_lyrics(slide, box_top_in, box_height_in, paragraphs, font_size):
    left = Inches(MARGIN_X_IN)
    top = Inches(box_top_in)
    width = Inches(SLIDE_WIDTH_IN - 2 * MARGIN_X_IN)
    height = Inches(box_height_in)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first_para_overall = True
    for para_lines in paragraphs:
        for i, line in enumerate(para_lines):
            if first_para_overall:
                p = tf.paragraphs[0]
                first_para_overall = False
            else:
                p = tf.add_paragraph()
            p.alignment = LYRIC_ALIGN
            p.line_spacing = LINE_SPACING_FACTOR
            # espaço extra só antes da 1a linha de cada estrofe (exceto a 1a estrofe do slide)
            if i == 0 and para_lines is not paragraphs[0]:
                p.space_before = Pt(PARAGRAPH_GAP_PT)
            run = p.add_run()
            run.text = line if line.strip() else " "
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.name = LYRIC_FONT_NAME
            run.font.color.rgb = LYRIC_COLOR


def build_pptx(slides: list, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    box_width_in = SLIDE_WIDTH_IN - 2 * MARGIN_X_IN
    lyrics_top_in = MARGIN_TOP_IN + TITLE_HEIGHT_IN
    lyrics_height_in = SLIDE_HEIGHT_IN - lyrics_top_in - MARGIN_BOTTOM_IN

    for song_slide in slides:
        plan = plan_slides_for_song(song_slide.paragraphs, box_width_in, lyrics_height_in)
        multi = len(plan) > 1

        for idx, (paras, font_size) in enumerate(plan, start=1):
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, prs)

            title_text = song_slide.title
            if multi:
                title_text = f"{song_slide.title} - Parte {idx}"
            add_title(slide, prs, title_text)

            add_lyrics(slide, lyrics_top_in, lyrics_height_in, paras, font_size)

    prs.save(output_path)


# ============================================================
# MAIN
# ============================================================


def main():
    parser = argparse.ArgumentParser(description="Gera PPTX 9:16 de letras de música a partir de um .txt")
    parser.add_argument("input", help="Arquivo .txt com '### Título' por slide")
    parser.add_argument("output", help="Caminho do .pptx de saída")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        text = f.read()

    slides = parse_input(text)
    if not slides:
        print("Nenhum slide encontrado. Confira se o .txt usa '### Título' antes de cada letra.", file=sys.stderr)
        sys.exit(1)

    build_pptx(slides, args.output)
    print(f"OK: {len(slides)} musica(s) processada(s) -> {args.output}")


if __name__ == "__main__":
    main()